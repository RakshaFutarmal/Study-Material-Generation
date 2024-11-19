from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PresentationGenerator:
    def __init__(self, txt_file_path, image_path):
        self.txt_file_path = txt_file_path
        self.image_path = image_path

    def parse_txt_file(self):
        slides = []

        with open(self.txt_file_path, 'r') as file:
            lines = file.readlines()

        current_slide = {"title": "", "content": []}
        for line in lines:
            line = line.strip()
            if line:
                if not current_slide["title"]:
                    current_slide["title"] = line
                else:
                    current_slide["content"].append(line)
            else:
                if current_slide["title"]:
                    slides.append(current_slide)
                    current_slide = {"title": "", "content": []}

        if current_slide["title"]:
            slides.append(current_slide)

        return slides

    def create_presentation(self):
        presentation = Presentation()

        first_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        img_left = Inches(0)
        img_top = Inches(0)
        img_width = Inches(10)
        img_height = Inches(7.5)
        first_slide.shapes.add_picture(self.image_path, img_left, img_top, img_width, img_height)

        slides_content = self.parse_txt_file()

        left_margin = Inches(0.5)
        right_margin = Inches(0.5)
        title_top_margin = Inches(0.5)
        content_top_margin = Inches(2.0)

        slide_width = Inches(10) - (left_margin + right_margin)

        for slide_content in slides_content:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

            title_box_height = Inches(1.0)
            title_box = slide.shapes.add_textbox(left_margin, title_top_margin, slide_width, title_box_height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            title_paragraph = title_frame.add_paragraph()
            title_paragraph.text = slide_content.get("title", "Slide Title")
            title_paragraph.font.size = Pt(24)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            title_paragraph.alignment = PP_ALIGN.CENTER

            content_box = slide.shapes.add_textbox(left_margin, content_top_margin, slide_width, Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for line in slide_content.get("content", []):
                content_paragraph = content_frame.add_paragraph()
                content_paragraph.text = line
                content_paragraph.font.size = Pt(16)
                content_paragraph.space_after = Pt(12)
                content_paragraph.alignment = PP_ALIGN.LEFT
                content_paragraph.font.color.rgb = RGBColor(0, 0, 0)

        file_name = "History of Aeronautics.pptx"
        presentation.save(file_name)
        return file_name

txt_filename = "/home/tntra/Information details/Study Material/Study Material using PPT/History of Aeronautics_content.txt"
image_path = "/home/tntra/Information details/Study Material/assests/Sample Tntra format.png"


presentation_generator = PresentationGenerator(txt_filename, image_path)
ppt_file = presentation_generator.create_presentation()

print(f"PowerPoint presentation generated and saved as '{ppt_file}'.")
















#######################################################################################################################################################

# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# from pptx.util import Inches
#
#
# def parse_txt_file(file_path):
#     slides = []
#     with open(file_path, 'r') as file:
#         lines = file.readlines()
#
#     current_slide = {"title": "", "content": []}
#     for line in lines:
#         line = line.strip()
#         if line:
#             if not current_slide["title"]:
#                 current_slide["title"] = line
#             else:
#                 current_slide["content"].append(line)
#         else:
#             if current_slide["title"]:
#                 slides.append(current_slide)
#                 current_slide = {"title": "", "content": []}
#
#     if current_slide["title"]:
#         slides.append(current_slide)
#
#     return slides
#
# def create_presentation_from_txt(file_path, image_path):
#     presentation = Presentation()
#
#     first_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
#     img_left = Inches(0)
#     img_top = Inches(0)
#     img_width = Inches(10)
#     img_hieght = Inches(7.5)
#     first_slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_hieght)
#
#     slides_content = parse_txt_file(file_path)
#
#     left_margin = Inches(0.5)
#     right_margin = Inches(0.5)
#     title_top_margin = Inches(0.5)
#     content_top_margin = Inches(2.0)
#
#     slide_width = Inches(10) - (left_margin + right_margin)
#
#     for slide_content in slides_content:
#         slide = presentation.slides.add_slide(presentation.slide_layouts[5])
#
#         background = slide.background
#         fill = background.fill
#         fill.solid()
#         fill.fore_color.rgb = RGBColor(255, 255, 255)
#
#         title_box_height = Inches(1.0)
#         title_box = slide.shapes.add_textbox(left_margin, title_top_margin, slide_width, title_box_height)
#         title_frame = title_box.text_frame
#         title_frame.word_wrap = True
#
#         title_paragraph = title_frame.add_paragraph()
#         title_paragraph.text = slide_content.get("title", "Slide Title")
#         title_paragraph.font.size = Pt(24)
#         title_paragraph.font.bold = True
#         title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
#         title_paragraph.alignment = PP_ALIGN.CENTER
#
#         content_box = slide.shapes.add_textbox(left_margin, content_top_margin, slide_width, Inches(5))
#         content_frame = content_box.text_frame
#         content_frame.word_wrap = True
#
#         for line in slide_content.get("content", []):
#             content_paragraph = content_frame.add_paragraph()
#             content_paragraph.text = line
#             content_paragraph.font.size = Pt(16)
#             content_paragraph.space_after = Pt(12)
#             content_paragraph.alignment = PP_ALIGN.LEFT
#             content_paragraph.font.color.rgb = RGBColor(0, 0, 0)
#
#     file_name = "History of Aeronautics.pptx"
#     presentation.save(file_name)
#     return file_name
#
#
# txt_filename = "/home/tntra/Information details/Study Material/Study Material using PPT/History of Aeronautics_content.txt"
# image_path = "/home/tntra/Documents/NeoTech/professional-mobility-test-develop/NeoTech Video Content/Sample Tntra format.png"
# ppt_file = create_presentation_from_txt(txt_filename, image_path)
# print(f"PowerPoint presentation generated and saved as '{ppt_file}'.")








