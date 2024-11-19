from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import integrated_function


class PresentationContentService:
    def content(self, subtopic, level):
        """
        Generate and add content for a specific subtopic and level.
        """
        query = f""" 
                    "Title": {subtopic},
                    "Table of content" : Generate 3-4 points for {subtopic},

                            "slide_1": Introduction to {subtopic},
                            "key_points": [
                                Provide a brief definition or introduction to {subtopic}.,
                                Explain why {subtopic} is important in the context of {level} audience.,
                                Mention any key applications or uses of {subtopic}.                                                  

                            "slide_2": Core Problem in {subtopic},
                            "key_points": [
                                Outline the main problem or challenge related to {subtopic}.,
                                Why is this problem significant for {level} audience?,
                                Provide examples of this challenge in real-world scenarios.

                            "slide_3": Importance of {subtopic},
                            "key_points": [
                                Highlight why {subtopic} is important for professionals or learners at the {level} level.,
                                Discuss its impact on the industry, technology, or society.,
                                Include statistics or research data if available.

                            "slide_4": Key Concepts in {subtopic},
                            "key_points": [
                                Define and explain 3-5 key concepts or terms within {subtopic}.,
                                Provide explanations at a {level}-appropriate level of detail.,
                                Use analogies, if necessary, for easier understanding at the {level} level.

                            slide_5": Proposed Solution for {subtopic} Challenges",
                            "key_points": [
                                Describe a known solution or approach to address challenges in {subtopic}.,
                                Explain why this solution is effective at the {level} of understanding.,
                                Provide any implementation examples or case studies.

                            "slide_6": Benefits of the Proposed Solution,
                            "key_points": [
                                Explain the benefits of the proposed solution in {subtopic}.,
                                Provide evidence of its effectiveness, tailored to the {level} audience.,
                                Compare this solution to alternative methods, if applicable.

                            "slide_7": Real-World Example of {subtopic},
                            "key_points": [
                                Present a real-world case study where {subtopic} was applied successfully.,
                                Discuss the challenges faced and how they were overcome.,
                                Include relevant metrics or outcomes to demonstrate impact.

                            "slide_8": Actionable Insights in {subtopic},
                            "key_points": [
                                Summarize the key takeaways for {level} learners or professionals.,
                                Provide practical steps they can apply to their work or studies.,
                                Discuss how {subtopic} can be integrated into real-life scenarios for their benefit.

                            "slide_9": Conclusion,
                            "key_points": [
                                Recap the most important points from the presentation.,
                                Tailor the conclusion to the {level} audience, focusing on what they need to retain.,
                                Encourage further exploration or study of {subtopic}.

                            "slide_10": Future Outlook for {subtopic},
                            "key_points": [
                                Provide predictions or recommendations for future trends in {subtopic}.,
                                Explain how the field may evolve, especially for the {level} audience.,
                                Highlight any opportunities for further research or innovation.

                        INSTRUCTIONS: The table of content should be a list based on the slide titles."""

        result1 = integrated_function.mistral_search(query)
        return result1

    def parse_generated_content(self, content):
        """
        Parse the generated content from the result.
        """
        slides = []
        lines = content.split('\n')

        current_slide = {"title": "", "content": []}
        for line in lines:
            line = line.strip()
            if line:
                if not current_slide["title"]:
                    current_slide["title"] = line
                else:
                    current_slide["content"].append(line)
            else:
                if current_slide["title"]:
                    slides.append(current_slide)
                    current_slide = {"title": "", "content": []}

        if current_slide["title"]:
            slides.append(current_slide)

        return slides

    def create_presentation_from_content(slides_content, image_path):
        """
        Create a PowerPoint presentation from the parsed content.
        """
        presentation = Presentation()

        first_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        img_left = Inches(0)
        img_top = Inches(0)
        img_width = Inches(10)
        img_height = Inches(7.5)
        first_slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)

        left_margin = Inches(0.5)
        right_margin = Inches(0.5)
        title_top_margin = Inches(0.5)
        content_top_margin = Inches(2.0)
        slide_width = Inches(10) - (left_margin + right_margin)

        for slide_content in slides_content:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

            title_box_height = Inches(1.0)
            title_box = slide.shapes.add_textbox(left_margin, title_top_margin, slide_width, title_box_height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            title_paragraph = title_frame.add_paragraph()
            title_paragraph.text = slide_content.get("title", "Slide Title")
            title_paragraph.font.size = Pt(24)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            title_paragraph.alignment = PP_ALIGN.CENTER

            content_box = slide.shapes.add_textbox(left_margin, content_top_margin, slide_width, Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for line in slide_content.get("content", []):
                content_paragraph = content_frame.add_paragraph()
                content_paragraph.text = line
                content_paragraph.font.size = Pt(16)
                content_paragraph.space_after = Pt(12)
                content_paragraph.alignment = PP_ALIGN.LEFT
                content_paragraph.font.color.rgb = RGBColor(0, 0, 0)

        file_name = "History_of_Aeronautics.pptx"
        presentation.save(file_name)
        return file_name

    def main():
        subtopics = input("Enter the subtopics: ").split(',')
        levels = input("Enter the levels: ").split(',')

        for subtopic, level in zip(subtopics, levels):
            service = PresentationContentService()
            content = service.content(subtopic, level)

            slides_content = service.parse_generated_content(content)

            image_path = "/home/tntra/Documents/NeoTech/professional-mobility-test-develop/NeoTech Video Content/Sample Tntra format.png"

            ppt_file = PresentationContentService.create_presentation_from_content(slides_content, image_path)

            print(f"PowerPoint presentation generated and saved as '{ppt_file}'.")


PresentationContentService.main()
